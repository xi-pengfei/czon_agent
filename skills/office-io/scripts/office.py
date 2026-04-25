#!/usr/bin/env python3
import argparse
import csv
import json
import mimetypes
import textwrap
from pathlib import Path


def emit(data):
    print(json.dumps(data, ensure_ascii=False, indent=2))


def fail(message, code=1):
    emit({"ok": False, "error": message})
    raise SystemExit(code)


def read_text(path: Path):
    return path.read_text(encoding="utf-8", errors="replace")


def read_csv(path: Path):
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        rows = list(csv.DictReader(f))
    return {"rows": rows, "count": len(rows)}


def read_docx(path: Path):
    try:
        from docx import Document
    except ImportError:
        fail("缺少依赖 python-docx，请先安装 requirements.txt")

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        tables.append(rows)
    return {"paragraphs": paragraphs, "tables": tables}


def read_xlsx(path: Path):
    try:
        from openpyxl import load_workbook
    except ImportError:
        fail("缺少依赖 openpyxl，请先安装 requirements.txt")

    wb = load_workbook(str(path), data_only=True, read_only=True)
    sheets = {}
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(list(row))
        sheets[ws.title] = rows
    return {"sheets": sheets}


def read_pdf(path: Path):
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            fail("缺少依赖 pypdf，请先安装 requirements.txt")

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        pages.append({"page": i, "text": page.extract_text() or ""})
    return {"pages": pages, "page_count": len(pages)}


def read_pptx(path: Path):
    try:
        from pptx import Presentation
    except ImportError:
        fail("缺少依赖 python-pptx，请先安装 requirements.txt")

    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        slides.append({"slide": index, "texts": texts})
    return {"slides": slides, "slide_count": len(slides)}


def inspect_file(path: Path):
    if not path.exists():
        fail(f"文件不存在：{path}")
    return {
        "path": str(path),
        "name": path.name,
        "suffix": path.suffix.lower(),
        "mime": mimetypes.guess_type(str(path))[0] or "application/octet-stream",
        "size": path.stat().st_size,
    }


def read_file(path: Path):
    info = inspect_file(path)
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt"}:
        payload = {"text": read_text(path)}
    elif suffix == ".csv":
        payload = read_csv(path)
    elif suffix == ".docx":
        payload = read_docx(path)
    elif suffix == ".xlsx":
        payload = read_xlsx(path)
    elif suffix == ".pdf":
        payload = read_pdf(path)
    elif suffix == ".pptx":
        payload = read_pptx(path)
    elif suffix in {".doc", ".xls", ".ppt"}:
        fail(f"暂不支持旧版二进制格式 {suffix}，请转换为新版格式后再处理")
    else:
        fail(f"暂不支持的文件类型：{suffix or '无后缀'}")
    return {"file": info, "content": payload}


def ensure_parent(path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)


def write_md(path: Path, text: str):
    ensure_parent(path)
    path.write_text(text, encoding="utf-8")
    return {"path": str(path), "bytes": path.stat().st_size}


def write_docx(path: Path, text: str):
    try:
        from docx import Document
    except ImportError:
        fail("缺少依赖 python-docx，请先安装 requirements.txt")

    ensure_parent(path)
    doc = Document()
    blocks = text.split("\n\n")
    for index, block in enumerate(blocks):
        block = block.strip()
        if not block:
            continue
        if index == 0 and len(block) <= 80:
            doc.add_heading(block, level=1)
        else:
            for line in block.splitlines():
                doc.add_paragraph(line)
    doc.save(str(path))
    return {"path": str(path), "bytes": path.stat().st_size}


def rows_from_json(raw: str):
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        fail(f"JSON 解析失败：{e}")
    if not isinstance(data, list):
        fail("JSON 必须是数组")
    return data


def write_xlsx(path: Path, raw_json: str):
    try:
        from openpyxl import Workbook
    except ImportError:
        fail("缺少依赖 openpyxl，请先安装 requirements.txt")

    rows = rows_from_json(raw_json)
    ensure_parent(path)
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    if rows and all(isinstance(row, dict) for row in rows):
        headers = list(rows[0].keys())
        ws.append(headers)
        for row in rows:
            ws.append([row.get(h) for h in headers])
    else:
        for row in rows:
            ws.append(row if isinstance(row, list) else [row])
    wb.save(str(path))
    return {"path": str(path), "rows": len(rows), "bytes": path.stat().st_size}


def write_pptx(path: Path, raw_json: str):
    try:
        from pptx import Presentation
    except ImportError:
        fail("缺少依赖 python-pptx，请先安装 requirements.txt")

    slides = rows_from_json(raw_json)
    ensure_parent(path)
    prs = Presentation()
    for item in slides:
        if not isinstance(item, dict):
            item = {"title": str(item), "bullets": []}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(item.get("title", ""))
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = item.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for index, bullet in enumerate(bullets):
            p = body.paragraphs[0] if index == 0 else body.add_paragraph()
            p.text = str(bullet)
            p.level = 0
    prs.save(str(path))
    return {"path": str(path), "slides": len(slides), "bytes": path.stat().st_size}


def write_pdf(path: Path, text: str):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.pdfgen import canvas
    except ImportError:
        fail("缺少依赖 reportlab，请先安装 requirements.txt")

    ensure_parent(path)
    c = canvas.Canvas(str(path), pagesize=A4)
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        c.setFont("STSong-Light", 11)
    except Exception:
        c.setFont("Helvetica", 11)
    width, height = A4
    x = 50
    y = height - 50
    for paragraph in text.splitlines() or [""]:
        lines = textwrap.wrap(paragraph, width=86) or [""]
        for line in lines:
            if y < 50:
                c.showPage()
                y = height - 50
            c.drawString(x, y, line)
            y -= 16
    c.save()
    return {"path": str(path), "bytes": path.stat().st_size}


def main():
    parser = argparse.ArgumentParser(description="Read and write common office files.")
    sub = parser.add_subparsers(dest="command", required=True)

    p = sub.add_parser("inspect")
    p.add_argument("path")

    p = sub.add_parser("read")
    p.add_argument("path")

    for name in ("write-md", "write-txt", "write-csv"):
        p = sub.add_parser(name)
        p.add_argument("path")
        p.add_argument("--text", required=True)

    p = sub.add_parser("write-docx")
    p.add_argument("path")
    p.add_argument("--text", required=True)

    p = sub.add_parser("write-xlsx")
    p.add_argument("path")
    p.add_argument("--json", required=True)

    p = sub.add_parser("write-pptx")
    p.add_argument("path")
    p.add_argument("--json", required=True)

    p = sub.add_parser("write-pdf")
    p.add_argument("path")
    p.add_argument("--text", required=True)

    args = parser.parse_args()
    path = Path(args.path)

    if args.command == "inspect":
        result = inspect_file(path)
    elif args.command == "read":
        result = read_file(path)
    elif args.command in {"write-md", "write-txt", "write-csv"}:
        result = write_md(path, args.text)
    elif args.command == "write-docx":
        result = write_docx(path, args.text)
    elif args.command == "write-xlsx":
        result = write_xlsx(path, args.json)
    elif args.command == "write-pptx":
        result = write_pptx(path, args.json)
    elif args.command == "write-pdf":
        result = write_pdf(path, args.text)
    else:
        fail(f"未知命令：{args.command}")

    emit({"ok": True, "data": result})


if __name__ == "__main__":
    main()
